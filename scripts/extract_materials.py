from __future__ import annotations

import hashlib
import json
import re
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[1]
SOURCE_DIR = ROOT / "考试素材合集"
OUTPUT_DIR = ROOT / "content" / "materials"
CATALOG_PATH = ROOT / "content" / "catalog.json"


@dataclass(frozen=True)
class Material:
    source: str
    slug: str
    category: str
    status: str | None = None


MATERIALS = (
    Material(
        '2026年“四个讲明”形势任务教育宣讲提纲（南货航）.doc',
        "four-explanations-cargo-2026",
        "形势任务",
    ),
    Material(
        '2026年“四个讲明”形势任务教育宣讲提纲（含集团和物流公司）.pptx',
        "four-explanations-group-logistics-2026",
        "形势任务",
    ),
    Material(
        "2026年南货航飞行队伍建设工作汇报材料暨飞行部年中工作报告 V3（无高亮，带批注）.doc",
        "flight-team-midyear-report-2026",
        "工作报告",
    ),
    Material("CCAR-398R1规章宣贯 202606V2.pdf", "ccar-398r1-2026", "安全规章"),
    Material("kaoti.doc", "exam-reference", "考试参考"),
    Material(
        '关于印发《2026年“四个讲明”形势任务教育宣讲提纲》的通知.pdf',
        "four-explanations-notice-2026",
        "形势任务",
    ),
    Material(
        "关于印发南航高质量发展总体思路（2025年版）的通知.pdf",
        "high-quality-development-2025",
        "战略发展",
        "旧材料",
    ),
    Material("南航物流改革发展.pptx", "logistics-reform", "改革发展"),
    Material("南货航2026年工作报告.docx", "cargo-report-2026", "工作报告"),
    Material(
        "南货航飞行部2025年工作报告 V7（终稿）.docx",
        "flight-report-2025",
        "工作报告",
        "旧材料",
    ),
    Material(
        '附件1：新时代“阳光南航”文化体系和品牌理念体系宣讲参考提纲.docx',
        "sunshine-csair-culture",
        "企业文化",
    ),
)


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as file:
        for chunk in iter(lambda: file.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def normalize_markdown(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{4,}", "\n\n\n", text)
    return text.strip()


def extract_word(path: Path) -> str:
    source = path
    temporary_directory: tempfile.TemporaryDirectory[str] | None = None
    try:
        if path.suffix.lower() == ".doc":
            temporary_directory = tempfile.TemporaryDirectory(prefix="777plus-word-")
            source = Path(temporary_directory.name) / f"{path.stem}.docx"
            subprocess.run(
                [
                    "powershell",
                    "-NoProfile",
                    "-ExecutionPolicy",
                    "Bypass",
                    "-File",
                    str(ROOT / "scripts" / "convert-word.ps1"),
                    "-Source",
                    str(path.resolve()),
                    "-Destination",
                    str(source.resolve()),
                ],
                check=True,
            )

        result = subprocess.run(
            [
                "pandoc",
                str(source),
                "--from=docx",
                "--to=gfm",
                "--wrap=none",
                "--track-changes=all",
                f"--lua-filter={ROOT / 'scripts' / 'drop-images.lua'}",
            ],
            check=True,
            capture_output=True,
            text=True,
            encoding="utf-8",
        )
        return normalize_markdown(result.stdout)
    finally:
        if temporary_directory is not None:
            temporary_directory.cleanup()


def iter_shape_text(shape) -> list[str]:
    blocks: list[str] = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in sorted(shape.shapes, key=lambda item: (item.top, item.left)):
            blocks.extend(iter_shape_text(child))
        return blocks

    if getattr(shape, "has_text_frame", False):
        paragraphs = []
        for paragraph in shape.text_frame.paragraphs:
            value = paragraph.text.strip()
            if value:
                paragraphs.append(value)
        if paragraphs:
            blocks.append("\n".join(paragraphs))

    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            cells = [cell.text.strip().replace("\n", " / ") for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            blocks.append("\n".join(rows))
    return blocks


def extract_presentation(path: Path) -> str:
    presentation = Presentation(path)
    sections = []
    for index, slide in enumerate(presentation.slides, start=1):
        blocks = []
        for shape in sorted(slide.shapes, key=lambda item: (item.top, item.left)):
            blocks.extend(iter_shape_text(shape))
        body = "\n\n".join(blocks).strip()
        sections.append(f"## 幻灯片 {index}\n\n{body}".rstrip())
    return normalize_markdown("\n\n".join(sections))


def extract_pdf(path: Path) -> str:
    document = fitz.open(path)
    try:
        sections = []
        for index, page in enumerate(document, start=1):
            text = page.get_text("text", sort=True).strip()
            sections.append(f"## 第 {index} 页\n\n{text}".rstrip())
        return normalize_markdown("\n\n".join(sections))
    finally:
        document.close()


def extract(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".doc", ".docx"}:
        return extract_word(path)
    if suffix == ".pptx":
        return extract_presentation(path)
    if suffix == ".pdf":
        return extract_pdf(path)
    raise ValueError(f"Unsupported material format: {path.name}")


def write_material(material: Material, source: Path, body: str) -> dict[str, object]:
    title = source.stem
    output = OUTPUT_DIR / f"{material.slug}.md"
    header = (
        "---\n"
        f'title: "{title.replace(chr(34), chr(92) + chr(34))}"\n'
        f'source: "{source.name.replace(chr(34), chr(92) + chr(34))}"\n'
        f'category: "{material.category}"\n'
        f'format: "{source.suffix.lower().lstrip(".")}"\n'
        f'sha256: "{sha256(source)}"\n'
        "---\n\n"
    )
    output.write_text(header + body + "\n", encoding="utf-8", newline="\n")
    item: dict[str, object] = {
        "id": material.slug,
        "title": title,
        "source": source.name,
        "category": material.category,
        "format": source.suffix.lower().lstrip("."),
        "path": f"materials/{output.name}",
        "sha256": sha256(source),
        "characters": len(body),
        "lines": body.count("\n") + 1 if body else 0,
    }
    if material.status:
        item["status"] = material.status
    return item


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    expected_sources = {material.source for material in MATERIALS}
    actual_sources = {path.name for path in SOURCE_DIR.iterdir() if path.is_file()}
    missing = sorted(expected_sources - actual_sources)
    if missing:
        raise SystemExit(
            json.dumps(
                {"missing_sources": missing},
                ensure_ascii=False,
                indent=2,
            )
        )

    catalog = []
    for material in MATERIALS:
        source = SOURCE_DIR / material.source
        body = extract(source)
        if not body:
            raise RuntimeError(f"No text extracted from {source.name}")
        catalog.append(write_material(material, source, body))
        print(f"{source.name}: {len(body)} characters")

    CATALOG_PATH.parent.mkdir(parents=True, exist_ok=True)
    CATALOG_PATH.write_text(
        json.dumps(catalog, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
        newline="\n",
    )
    print(f"Wrote {len(catalog)} materials to {OUTPUT_DIR.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
